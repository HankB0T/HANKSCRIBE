#!/usr/bin/env python3
"""
Context Builder (Hankscribe 2.0)
Scans your project folder (from config.json) and builds a searchable JSON
index for fast Q&A retrieval. Nothing here is project-specific — point
config.json's paths.project_dir at your own folder and run this.
"""

import os
import json
import hashlib
import re
from pathlib import Path
from datetime import datetime
from collections import Counter


BASE_DIR = os.path.dirname(os.path.abspath(__file__))


def _load_paths():
    """Resolve the ACTIVE project's folder list + index filename. A project
    may span several folders (paths.project_dirs is a list), and config.json
    may define multiple named projects under "projects" with an "active" key."""
    prof = {}
    try:
        with open(os.path.join(BASE_DIR, "config.json"), "r", encoding="utf-8") as f:
            cfg = json.load(f)
        prof = dict(cfg.get("paths", {}))
        projects = cfg.get("projects") or {}
        active = projects.get("active")
        if active and isinstance(projects.get(active), dict):
            prof.update(projects[active])
            print(f"  Project: {active}")
    except FileNotFoundError:
        print("  ⚠ config.json not found — using default project dir")
    except Exception as e:
        print(f"  ⚠ config.json unreadable ({e}) — using defaults")
    dirs = prof.get("project_dirs") or prof.get("project_dir") or "~/Desktop/MyProject"
    if isinstance(dirs, str):
        dirs = [dirs]
    index_name = prof.get("context_index", "context-index.json")
    return [os.path.expanduser(d) for d in dirs], os.path.join(BASE_DIR, index_name)


PROJECT_DIRS, OUTPUT_FILE = _load_paths()

SUPPORTED_EXTENSIONS = {'.md', '.txt', '.vtt', '.docx', '.pdf', '.pptx', '.olm'}


def extract_text(filepath):
    ext = filepath.suffix.lower()
    handlers = {
        '.md': extract_plain,
        '.txt': extract_plain,
        '.vtt': extract_vtt,
        '.docx': extract_docx,
        '.pdf': extract_pdf,
        '.pptx': extract_pptx,
        '.olm': extract_olm,
    }
    handler = handlers.get(ext)
    if handler:
        try:
            return handler(filepath)
        except Exception as e:
            return f"[Error: {e}]"
    return None


def extract_plain(filepath):
    with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
        return f.read().strip()


def extract_vtt(filepath):
    with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
        content = f.read()
    content = re.sub(r'^WEBVTT.*?\n\n', '', content, flags=re.MULTILINE)
    content = re.sub(r'\d{2}:\d{2}:\d{2}\.\d{3}\s*-->\s*\d{2}:\d{2}:\d{2}\.\d{3}', '', content)
    content = re.sub(r'[0-9a-f]{8}-[0-9a-f]{4}-[0-9a-f]{4}-[0-9a-f]{4}-[0-9a-f]{12}-\d+', '', content)
    content = re.sub(r'^\d+$', '', content, flags=re.MULTILINE)
    lines = [line.strip() for line in content.split('\n') if line.strip()]
    return ' '.join(lines)


def extract_docx(filepath):
    try:
        from docx import Document
        doc = Document(filepath)
        return '\n'.join(p.text.strip() for p in doc.paragraphs if p.text.strip())
    except ImportError:
        return "[python-docx not installed]"


def extract_pdf(filepath):
    try:
        import pdfplumber
        text = []
        with pdfplumber.open(filepath) as pdf:
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    text.append(t.strip())
        return '\n\n'.join(text)
    except ImportError:
        return "[pdfplumber not installed]"


def extract_pptx(filepath):
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        text = []
        for i, slide in enumerate(prs.slides, 1):
            parts = [f"[Slide {i}]"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
            if len(parts) > 1:
                text.append('\n'.join(parts))
        return '\n\n'.join(text)
    except ImportError:
        return "[python-pptx not installed]"


def extract_olm(filepath):
    try:
        import zipfile
        text = []
        with zipfile.ZipFile(filepath, 'r') as z:
            for name in z.namelist():
                if name.endswith('.xml') or name.endswith('.eml'):
                    try:
                        content = z.read(name).decode('utf-8', errors='ignore')
                        content = re.sub(r'<[^>]+>', '', content)
                        content = ' '.join(content.split())
                        if content:
                            text.append(content)
                    except:
                        pass
        return '\n\n'.join(text) if text else "[Empty OLM]"
    except Exception as e:
        return f"[OLM error: {e}]"


def extract_keywords(text):
    words = re.findall(r'\b[A-Z][a-z]+(?:\s+[A-Z][a-z]+)*\b', text)
    common = {'The', 'This', 'That', 'These', 'Those', 'There', 'When', 'Where', 'What', 'Does'}
    keywords = [w for w in words if w not in common and len(w) > 3]
    return [k for k, _ in Counter(keywords).most_common(20)]


def _load_previous_index():
    """Previous index keyed by relative path, for incremental rebuilds."""
    try:
        with open(OUTPUT_FILE, 'r', encoding='utf-8') as f:
            old = json.load(f)
        return {fi['path']: fi for fi in old.get('files', [])}
    except Exception:
        return {}


def build_context():
    roots = [Path(d) for d in PROJECT_DIRS if Path(d).exists()]
    missing = [d for d in PROJECT_DIRS if not Path(d).exists()]
    for m in missing:
        print(f"  ⚠ Folder not found (skipped): {m}")
    if not roots:
        print(f"  ✗ No project folder exists — set project_dirs in config.json")
        return None

    previous = _load_previous_index()
    reused = 0

    context_data = {
        'built_at': datetime.now().isoformat(),
        'source_dirs': [str(r) for r in roots],
        'files': [],
        'file_count': 0,
        'total_chars': 0
    }

    all_files = []
    for root in roots:
        print(f"  Scanning: {root}")
        for filepath in sorted(root.rglob('*')):
            if filepath.is_file() and filepath.suffix.lower() in SUPPORTED_EXTENSIONS:
                all_files.append((root, filepath))

    for project_path, filepath in all_files:
        if True:
            # Skip huge files (>500KB text) — likely raw email dumps
            if filepath.stat().st_size > 500_000 and filepath.suffix.lower() == '.olm':
                print(f"    Skip (too large): {filepath.name}")
                continue

            # Prefix with the root folder's name so files from different
            # folders can't collide in the index
            rel_path = f"{project_path.name}/{filepath.relative_to(project_path)}"
            stat = filepath.stat()
            modified = datetime.fromtimestamp(stat.st_mtime).isoformat()

            # Incremental: reuse the previous extraction when the file hasn't
            # changed (same mtime + size) — text extraction, especially for
            # PDFs/DOCX, dominates rebuild time.
            prev = previous.get(rel_path)
            if prev and prev.get('modified') == modified and prev.get('size') == stat.st_size:
                context_data['files'].append(prev)
                context_data['file_count'] += 1
                context_data['total_chars'] += prev.get('char_count', len(prev.get('text', '')))
                reused += 1
                continue

            text = extract_text(filepath)
            if text and len(text) > 10:
                file_info = {
                    'name': filepath.name,
                    'path': rel_path,
                    'type': filepath.suffix.lower(),
                    'size': stat.st_size,
                    'modified': modified,
                    'hash': hashlib.md5(filepath.read_bytes()).hexdigest(),
                    'char_count': len(text),
                    'keywords': extract_keywords(text)[:10],
                    'text': text
                }
                context_data['files'].append(file_info)
                context_data['file_count'] += 1
                context_data['total_chars'] += len(text)
                print(f"    ✓ {filepath.name} ({len(text):,} chars)")

    if reused:
        print(f"    ↻ {reused} unchanged file(s) reused from previous index")
    return context_data


if __name__ == '__main__':
    print("\n  Building project context for Hankscribe 2.0...")
    print("  " + "=" * 50)

    context = build_context()
    if context:
        with open(OUTPUT_FILE, 'w', encoding='utf-8') as f:
            json.dump(context, f, ensure_ascii=False, indent=2)
        print(f"\n  ✓ Done: {context['file_count']} files, {context['total_chars']:,} chars")
        print(f"  Saved: {OUTPUT_FILE}\n")
    else:
        print("\n  ✗ Failed to build context\n")
